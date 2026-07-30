"""Auditable metadata retrieval over the local licensed 5·3 question corpus.

Raw copyrighted files stay outside Git.  The committed catalog contains only
source metadata.  Short content previews are extracted on demand and never
included in student-facing API responses.
"""

from __future__ import annotations

import hashlib
import json
import re
from collections import Counter
from collections.abc import Iterable
from pathlib import Path
from typing import Any

from ai_education.domain.enums import Subject
from ai_education.domain.homework import QuestionBankEvidence

PROJECT_ROOT = Path(__file__).resolve().parents[3]
CORPUS_ROOT = PROJECT_ROOT / "Knowledge" / "title" / "2026五年高考三年模拟53A、B版新高考全套资料"
CATALOG_PATH = PROJECT_ROOT / "Knowledge" / "catalogs" / "question_bank_catalog.json"
INDEX_PATH = PROJECT_ROOT / "Knowledge" / "91_indexes" / "question_bank.db"

SUBJECT_MARKERS: tuple[tuple[str, Subject], ...] = (
    ("思想政治", Subject.IDEOLOGY_POLITICS),
    ("政治", Subject.IDEOLOGY_POLITICS),
    ("语文", Subject.CHINESE),
    ("数学", Subject.MATHEMATICS),
    ("英语", Subject.FOREIGN_LANGUAGE),
    ("物理", Subject.PHYSICS),
    ("化学", Subject.CHEMISTRY),
    ("生物", Subject.BIOLOGY),
    ("历史", Subject.HISTORY),
    ("地理", Subject.GEOGRAPHY),
    ("信息技术", Subject.TECHNOLOGY),
    ("通用技术", Subject.TECHNOLOGY),
    ("技术", Subject.TECHNOLOGY),
)

REGION_MARKERS = (
    "全国新高考",
    "新高考",
    "北京",
    "天津",
    "山东",
    "广东",
    "江苏",
    "浙江",
    "湖南",
    "西北",
)

SUPPORTED_SUFFIXES = {".pdf", ".docx", ".pptx", ".pptm", ".txt", ".mp3"}


def source_id(relative_path: str) -> str:
    return f"qb_{hashlib.sha256(relative_path.encode('utf-8')).hexdigest()[:18]}"


def classify_path(relative_path: str, size: int = 0) -> dict[str, Any]:
    normalized = relative_path.replace("\\", "/")
    path = Path(normalized)
    joined = "/".join(path.parts)
    subject = next((item.value for marker, item in SUBJECT_MARKERS if marker in joined), None)
    edition = (
        "A"
        if normalized.startswith("A版/")
        else "B"
        if normalized.startswith("B版/")
        else "unknown"
    )
    region = next((marker for marker in REGION_MARKERS if f"{marker}版" in joined), None)
    if region is None:
        region = "全国新高考" if "新高考" in joined or edition == "A" else "地区待确认"

    if any(marker in joined for marker in ("答案", "解析", "精析")):
        role = "answer_secure"
    elif any(marker in joined for marker in ("讲解册", "精讲", "知识清单", "技巧清单")):
        role = "explanation_secure"
    elif any(marker in joined for marker in ("训练", "精练", "真题", "预测", "题型", "试题")):
        role = "exercise"
    elif "封面" in joined:
        role = "cover"
    elif path.suffix.lower() == ".mp3":
        role = "audio"
    elif path.suffix.lower() in {".zip", ".rar"}:
        role = "archive"
    else:
        role = "supporting_material"

    topic = None
    for part in reversed(path.parts[:-1]):
        clean = re.sub(r"^[0-9_：:.、\s]+", "", part).strip()
        if any(marker in clean for marker in ("专题", "模块", "预测", "真题", "微专题")):
            topic = clean[:160]
            break
    if topic is None and len(path.parts) >= 2:
        topic = path.parts[-2][:160]

    return {
        "source_id": source_id(normalized),
        "relative_path": normalized,
        "title": path.name.strip() or path.name,
        "subject": subject,
        "edition": edition,
        "region": region,
        "content_role": role,
        "topic": topic,
        "file_type": path.suffix.lower().lstrip(".") or "unknown",
        "file_size": max(size, 0),
        "secure_content_available": role in {"answer_secure", "explanation_secure"},
    }


def build_catalog_items(root: Path = CORPUS_ROOT) -> list[dict[str, Any]]:
    root = root.resolve()
    if not root.exists():
        return []
    items = []
    for path in sorted(root.rglob("*")):
        if not path.is_file():
            continue
        relative = path.relative_to(root).as_posix()
        item = classify_path(relative, path.stat().st_size)
        item["supported"] = path.suffix.lower() in SUPPORTED_SUFFIXES
        items.append(item)
    return items


def catalog_summary(items: Iterable[dict[str, Any]]) -> dict[str, Any]:
    materialized = list(items)
    return {
        "total_files": len(materialized),
        "subjects": dict(
            sorted(Counter(item.get("subject") or "unclassified" for item in materialized).items())
        ),
        "editions": dict(sorted(Counter(item["edition"] for item in materialized).items())),
        "regions": dict(sorted(Counter(item["region"] for item in materialized).items())),
        "content_roles": dict(
            sorted(Counter(item["content_role"] for item in materialized).items())
        ),
        "file_types": dict(sorted(Counter(item["file_type"] for item in materialized).items())),
        "total_bytes": sum(int(item.get("file_size", 0)) for item in materialized),
    }


class QuestionBankService:
    """Retrieve source-grounded practice metadata without exposing answer files."""

    def __init__(self, catalog_path: Path = CATALOG_PATH, corpus_root: Path = CORPUS_ROOT) -> None:
        self.catalog_path = catalog_path
        self.corpus_root = corpus_root.resolve()
        self._items: list[dict[str, Any]] | None = None
        self._by_id: dict[str, dict[str, Any]] = {}

    def _load(self) -> list[dict[str, Any]]:
        if self._items is not None:
            return self._items
        if self.catalog_path.exists():
            payload = json.loads(self.catalog_path.read_text(encoding="utf-8"))
            self._items = list(payload.get("items", []))
        else:
            self._items = build_catalog_items(self.corpus_root)
        self._by_id = {item["source_id"]: item for item in self._items}
        return self._items

    def summary(self) -> dict[str, Any]:
        if self.catalog_path.exists():
            payload = json.loads(self.catalog_path.read_text(encoding="utf-8"))
            return payload.get("summary", catalog_summary(payload.get("items", [])))
        return catalog_summary(self._load())

    def search(
        self,
        query: str,
        *,
        subject: Subject | None = None,
        province: str | None = None,
        include_secure: bool = False,
        limit: int = 5,
    ) -> list[QuestionBankEvidence]:
        query = query.strip()
        tokens = self._query_tokens(query)
        candidates: list[tuple[float, dict[str, Any]]] = []
        for item in self._load():
            if not item.get("supported", True):
                continue
            if subject and item.get("subject") not in {subject.value, None}:
                continue
            if not include_secure and item.get("content_role") in {
                "answer_secure",
                "explanation_secure",
            }:
                continue
            haystack = " ".join(
                str(item.get(key, ""))
                for key in ("relative_path", "title", "topic", "region", "content_role")
            ).lower()
            score = 2.5 if subject and item.get("subject") == subject.value else 0.0
            score += 1.5 if item.get("content_role") == "exercise" else 0.2
            if province and province in str(item.get("region", "")):
                score += 1.2
            for token in tokens:
                if token.lower() in haystack:
                    score += min(2.0, 0.45 + len(token) * 0.08)
            if tokens and score <= 2.7:
                continue
            candidates.append((score, item))
        candidates.sort(key=lambda pair: (-pair[0], pair[1]["relative_path"]))
        results = []
        for score, item in candidates[: max(1, min(limit, 20))]:
            evidence_payload = {key: value for key, value in item.items() if key != "supported"}
            results.append(
                QuestionBankEvidence.model_validate(
                    {
                        **evidence_payload,
                        "confidence": round(min(0.98, 0.55 + score / 16), 3),
                    }
                )
            )
        return results

    def secure_preview(self, source: str, *, max_chars: int = 1200) -> str:
        """Extract a bounded internal preview; callers must never return it to students."""

        item = self._by_id.get(source) or next(
            (candidate for candidate in self._load() if candidate["source_id"] == source), None
        )
        if not item:
            return ""
        path = (self.corpus_root / item["relative_path"]).resolve()
        if self.corpus_root not in path.parents or not path.is_file():
            return ""
        try:
            suffix = path.suffix.lower()
            if suffix == ".txt":
                text = path.read_text(encoding="utf-8", errors="ignore")
            elif suffix == ".docx":
                from docx import Document

                text = "\n".join(paragraph.text for paragraph in Document(path).paragraphs)
            elif suffix in {".pptx", ".pptm"}:
                from pptx import Presentation

                text = "\n".join(
                    shape.text
                    for slide in Presentation(path).slides
                    for shape in slide.shapes
                    if hasattr(shape, "text") and shape.text
                )
            elif suffix == ".pdf" and path.stat().st_size <= 80 * 1024 * 1024:
                from pypdf import PdfReader

                reader = PdfReader(path)
                text = "\n".join((page.extract_text() or "") for page in reader.pages[:6])
            else:
                return ""
        except Exception:
            return ""
        return re.sub(r"\s+", " ", text).strip()[: max(100, min(max_chars, 4000))]

    @staticmethod
    def _query_tokens(query: str) -> list[str]:
        chunks = re.findall(r"[\u4e00-\u9fff]{2,10}|[A-Za-z0-9_+-]{2,30}", query)
        ignored = {"请问", "怎么", "如何", "这个", "题目", "作业", "答案", "学生"}
        unique: list[str] = []
        for chunk in chunks:
            if chunk in ignored or chunk in unique:
                continue
            unique.append(chunk)
        return unique[:12]
